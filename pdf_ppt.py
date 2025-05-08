import os
import tempfile
import threading
import tkinter as tk
from pathlib import Path
from pptx import Presentation
from db_connector import DBConnector
from pdf2image import convert_from_path
from tkinter import filedialog, messagebox, ttk

db = DBConnector()

class PowerPointConverterTab:
    def __init__(self, notebook, username):
        self.username = username
        self.tab = ttk.Frame(notebook)
        notebook.add(self.tab, text="PPT ↔ PDF")
        
        self.conversion_direction = tk.StringVar(value="pdf_to_ppt")
        self.conversion_direction.trace('w', self.on_conversion_direction_change)
        
        main_frame = ttk.Frame(self.tab)
        main_frame.pack(fill="both", expand=True, padx=10, pady=10)
        
        type_frame = ttk.LabelFrame(main_frame, text="Conversion Type")
        type_frame.pack(fill="x", pady=(0, 10))
        
        ttk.Radiobutton(type_frame, text="PPT/PPTX to PDF", 
                       variable=self.conversion_direction, 
                       value="ppt_to_pdf").pack(side="left", padx=5)
        
        ttk.Radiobutton(type_frame, text="PDF to PPT/PPTX (Images)", 
                       variable=self.conversion_direction, 
                       value="pdf_to_ppt").pack(side="left", padx=5)
        
        self.content_container = ttk.Frame(main_frame)
        self.content_container.pack(fill="both", expand=True)
        
        self.under_construction_frame = ttk.Frame(self.content_container)
        ttk.Label(self.under_construction_frame, 
                 text="PPT to PDF conversion is under maintenance",
                 font=('Helvetica', 12)).pack(pady=50)
        
        self.pdf_to_ppt_frame = ttk.Frame(self.content_container)
        
        file_frame = ttk.LabelFrame(self.pdf_to_ppt_frame, text="Files")
        file_frame.pack(fill="both", expand=True)
        
        self.file_listbox = tk.Listbox(file_frame, selectmode=tk.EXTENDED)
        self.file_listbox.pack(side="left", fill="both", expand=True, padx=5, pady=5)
        
        scrollbar = ttk.Scrollbar(file_frame, command=self.file_listbox.yview)
        scrollbar.pack(side="right", fill="y")
        self.file_listbox.config(yscrollcommand=scrollbar.set)
        
        button_frame = ttk.Frame(self.pdf_to_ppt_frame)
        button_frame.pack(fill="x", pady=(10, 0))
        
        ttk.Button(button_frame, text="Add Files", command=self.add_files).pack(side="left", padx=5)
        ttk.Button(button_frame, text="Remove Selected", command=self.remove_files).pack(side="left", padx=5)
        ttk.Button(button_frame, text="Clear All", command=self.clear_files).pack(side="left", padx=5)
        
        output_frame = ttk.Frame(self.pdf_to_ppt_frame)
        output_frame.pack(fill="x", pady=(10, 0))
        
        ttk.Label(output_frame, text="Output Folder:").pack(side="left")
        self.output_entry = ttk.Entry(output_frame)
        self.output_entry.pack(side="left", fill="x", expand=True, padx=5)
        ttk.Button(output_frame, text="Browse", command=self.browse_output).pack(side="left")
        
        ttk.Button(self.pdf_to_ppt_frame, text="Convert", command=self.start_conversion, 
                  style="Accent.TButton").pack(pady=(10, 0), fill="x")
        
        self.status_var = tk.StringVar(value="Ready")
        ttk.Label(self.pdf_to_ppt_frame, textvariable=self.status_var, relief=tk.SUNKEN,
                 anchor=tk.W).pack(fill="x", padx=10, pady=(0, 10))
        
        self.files = []
        self.on_conversion_direction_change()
        self.check_dependencies()
    
    def on_conversion_direction_change(self, *args):
        for widget in self.content_container.winfo_children():
            widget.pack_forget()
            
        if self.conversion_direction.get() == "ppt_to_pdf":
            self.under_construction_frame.pack(fill="both", expand=True)
        else:
            self.pdf_to_ppt_frame.pack(fill="both", expand=True)
    
    def check_dependencies(self):
        try:
            self.has_pdf2image = True
        except ImportError:
            self.has_pdf2image = False
            messagebox.showwarning("Requirements",
                "For PDF to PPT conversion, please install:\n"
                "pip install pdf2image pillow")
    
    def add_files(self):
        if self.conversion_direction.get() != "pdf_to_ppt":
            return
            
        filetypes = [("PDF Files", "*.pdf")]
        title = "Select PDF Files"
            
        files = filedialog.askopenfilenames(title=title, filetypes=filetypes)
        if files:
            self.files.extend(files)
            self.update_listbox()
            self.status_var.set(f"{len(self.files)} file(s) selected")
    
    def remove_files(self):
        selected = self.file_listbox.curselection()
        if selected:
            for i in reversed(selected):
                self.files.pop(i)
            self.update_listbox()
            self.status_var.set(f"{len(self.files)} file(s) remaining")
    
    def clear_files(self):
        self.files = []
        self.update_listbox()
        self.status_var.set("File list cleared")
    
    def browse_output(self):
        folder = filedialog.askdirectory(title="Select Output Folder")
        if folder:
            self.output_entry.delete(0, tk.END)
            self.output_entry.insert(0, folder)
    
    def update_listbox(self):
        self.file_listbox.delete(0, tk.END)
        for file in self.files:
            self.file_listbox.insert(tk.END, os.path.basename(file))
    
    def start_conversion(self):
        if self.conversion_direction.get() == "ppt_to_pdf":
            messagebox.showinfo("Information", "PPT to PDF conversion is currently under maintenance")
            return
            
        if not self.files:
            messagebox.showwarning("Warning", "No files selected!")
            return
            
        output_folder = self.output_entry.get()
        if not output_folder:
            messagebox.showwarning("Warning", "Please select an output folder!")
            return
            
        thread = threading.Thread(
            target=self.convert_pdf_to_ppt,
            args=(output_folder,),
            daemon=True
        )
        thread.start()
    
    def convert_pdf_to_ppt(self, output_folder):
        if not self.has_pdf2image:
            messagebox.showerror("Error", "pdf2image not installed. Cannot convert PDF to PPT.")
            return 0        
        success_count = 0
        
        for i, input_file in enumerate(self.files):
            try:
                self.status_var.set(f"Converting {i+1}/{len(self.files)}: {os.path.basename(input_file)}")
                
                output_file = os.path.join(output_folder, f"{Path(input_file).stem}.pptx")
                
                images = convert_from_path(input_file)
                prs = Presentation()
                
                for img in images:
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_img:
                        img.save(temp_img.name, "PNG")
                    
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    slide.shapes.add_picture(
                        temp_img.name, 
                        left=0, 
                        top=0, 
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                    
                    os.unlink(temp_img.name)
                
                prs.save(output_file)
                db.log_task("pdf_to_ppt", input_file, output_file)
                success_count += 1
                
            except Exception as e:
                print(f"Error converting {input_file}: {str(e)}")
        
        self.status_var.set(f"Conversion complete! {success_count}/{len(self.files)} files converted")
        messagebox.showinfo("Success", f"Converted {success_count} file(s)")
        return success_count

if __name__ == "__main__":
    root = tk.Tk()
    root.title("PPT/PDF Converter")
    
    notebook = ttk.Notebook(root)
    notebook.pack(fill="both", expand=True)
    
    PowerPointConverterTab(notebook)
    root.mainloop()